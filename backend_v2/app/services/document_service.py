try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None
import os
import docx
import openpyxl
from pptx import Presentation
from sqlalchemy.orm import Session
from app.models.ai import KnowledgeSource, KnowledgeChunk
from app.services.document_agent import clean_document_text

class DocumentProcessor:
    """
    Service to process academic materials (PDFs, PPTX, DOCX, XLSX) into indexed chunks for RAG.
    """
    
    @staticmethod
    def extract_text(file_path: str, filename: str = None) -> list:
        """
        Extracts raw text from a file based on extension.
        Returns a list of dicts: [{'page_number': int, 'content': str}]
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
            
        ext = os.path.splitext(filename or file_path)[1].lower()
        
        if ext == '.pdf':
            return DocumentProcessor._extract_pdf(file_path)
        elif ext in ['.docx', '.doc']:
            return DocumentProcessor._extract_docx(file_path)
        elif ext in ['.pptx', '.ppt']:
            return DocumentProcessor._extract_pptx(file_path)
        elif ext in ['.xlsx', '.xls']:
            return DocumentProcessor._extract_xlsx(file_path)
        elif ext in ['.txt', '.csv']:
            with open(file_path, 'r', encoding='utf-8') as f:
                return [{"page_number": 1, "content": clean_document_text(f.read())}]
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    @staticmethod
    def _extract_pdf(file_path: str):
        if fitz is None:
            raise ImportError("PyMuPDF (fitz) is not installed. PDF extraction is unavailable.")
        doc = fitz.open(file_path)
        pages_text = []
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text("text")
            pages_text.append({
                "page_number": page_num + 1,
                "content": clean_document_text(text)
            })
        doc.close()
        return pages_text

    @staticmethod
    def _extract_docx(file_path: str):
        doc = docx.Document(file_path)
        full_text = "\n".join([para.text for para in doc.paragraphs])
        return [{"page_number": 1, "content": clean_document_text(full_text)}]

    @staticmethod
    def _extract_pptx(file_path: str):
        prs = Presentation(file_path)
        pages_text = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            pages_text.append({
                "page_number": i + 1,
                "content": clean_document_text("\n".join(slide_text))
            })
        return pages_text

    @staticmethod
    def _extract_xlsx(file_path: str):
        wb = openpyxl.load_workbook(file_path, data_only=True)
        pages_text = []
        for i, sheet_name in enumerate(wb.sheetnames):
            sheet = wb[sheet_name]
            rows_text = []
            for row in sheet.iter_rows(values_only=True):
                rows_text.append(" ".join([str(cell) for cell in row if cell is not None]))
            pages_text.append({
                "page_number": i + 1,
                "content": clean_document_text("\n".join(rows_text))
            })
        return pages_text

    @staticmethod
    def chunk_text(pages_text, chunk_size=1200, overlap=200):
        """
        Splits extracted text into overlapping chunks.
        Keeps track of the page number.
        """
        chunks = []
        current_chunk = ""
        current_pages = []
        
        for page in pages_text:
            content = page["content"]
            page_num = page["page_number"]
            
            # Simple word-based chunking
            words = content.split()
            for word in words:
                if len(current_chunk) + len(word) + 1 > chunk_size:
                    # Save chunk
                    chunks.append({
                        "content": current_chunk.strip(),
                        "page_number": current_pages[0] if current_pages else page_num
                    })
                    # Start new chunk with overlap
                    current_chunk = current_chunk[-overlap:] if len(current_chunk) > overlap else current_chunk
                    current_chunk += " " + word
                    current_pages = [page_num]
                else:
                    current_chunk += " " + word
                    if page_num not in current_pages:
                        current_pages.append(page_num)
                        
        # Add final chunk
        if current_chunk:
            chunks.append({
                "content": current_chunk.strip(),
                "page_number": current_pages[0] if current_pages else 1
            })
            
        return chunks

    @classmethod
    def index_document(cls, db: Session, source_id: str):
        """
        Orchestrates the full indexing pipeline for a KnowledgeSource.
        """
        source = db.query(KnowledgeSource).filter(KnowledgeSource.id == source_id).first()
        if not source:
            print(f"Indexing Error: Source {source_id} not found.")
            return False
            
        try:
            source.status = "processing"
            db.commit()
            
            print(f"Indexing Started: {source.title} ({source.file_path})")
            
            # 1. Extract Text
            pages_text = cls.extract_text(source.file_path)
            source.page_count = len(pages_text)
            
            # 2. Chunk Text
            chunks_data = cls.chunk_text(pages_text)
            source.chunk_count = len(chunks_data)
            
            # 3. Save Chunks
            # Clear existing chunks if re-indexing
            db.query(KnowledgeChunk).filter(KnowledgeChunk.source_id == source_id).delete()
            
            for idx, chunk in enumerate(chunks_data):
                new_chunk = KnowledgeChunk(
                    source_id=source_id,
                    content=chunk["content"],
                    page_number=chunk["page_number"],
                    chunk_index=idx
                )
                db.add(new_chunk)
                
            source.status = "indexed"
            db.commit()
            print(f"Indexing Success: {source.title} - {len(chunks_data)} chunks created.")
            return True
            
        except Exception as e:
            print(f"Indexing Failed for source {source_id}: {str(e)}")
            source.status = "failed"
            db.commit()
            return False

document_service = DocumentProcessor()
